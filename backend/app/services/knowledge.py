"""Knowledge base ingestion + search.

Ingestion pipeline:
    text | url | file → fetch (trafilatura for URLs, extractor for files) →
    chunk → embed → store.

Search:
    embed(query) → pgvector cosine similarity → top_k chunks.
"""

from __future__ import annotations

import io
import logging
import uuid
from dataclasses import dataclass

import httpx
from sqlalchemy import desc, func, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.agents.harness.embedder import embed as embed_text
from app.core.errors import NotFound
from app.db.models.attachment import Attachment, AttachmentKind
from app.db.models.knowledge import (
    DocSourceKind,
    DocStatus,
    KnowledgeChunk,
    KnowledgeCollection,
    KnowledgeDoc,
)
from app.db.repository import AsyncRepository

log = logging.getLogger(__name__)


DEFAULT_CHUNK_SIZE = 800  # target chars per chunk
DEFAULT_CHUNK_OVERLAP = 80


# ─── Collection CRUD ────────────────────────────────────────
async def list_collections(
    session: AsyncSession, *, workspace_id: uuid.UUID
) -> list[tuple[KnowledgeCollection, int, int]]:
    """Return (collection, doc_count, chunk_count) for each collection."""
    stmt = (
        select(
            KnowledgeCollection,
            func.count(KnowledgeDoc.id.distinct()).label("doc_count"),
            func.coalesce(func.sum(KnowledgeDoc.chunk_count), 0).label("chunk_count"),
        )
        .outerjoin(
            KnowledgeDoc,
            (KnowledgeDoc.collection_id == KnowledgeCollection.id)
            & (KnowledgeDoc.deleted_at.is_(None)),
        )
        .where(
            KnowledgeCollection.workspace_id == workspace_id,
            KnowledgeCollection.deleted_at.is_(None),
        )
        .group_by(KnowledgeCollection.id)
        .order_by(desc(KnowledgeCollection.created_at))
    )
    rows = (await session.execute(stmt)).all()
    return [(r[0], int(r[1] or 0), int(r[2] or 0)) for r in rows]


async def get_collection_or_404(
    session: AsyncSession, collection_id: uuid.UUID, *, workspace_id: uuid.UUID
) -> KnowledgeCollection:
    repo = AsyncRepository(session, KnowledgeCollection)
    row = await repo.get(collection_id)
    if row is None or row.workspace_id != workspace_id or row.deleted_at is not None:
        raise NotFound("collection_not_found", code="knowledge.collection_not_found")
    return row


# ─── Document ingestion ────────────────────────────────────
@dataclass
class IngestResult:
    doc: KnowledgeDoc
    chunks: int


async def ingest_document(
    session: AsyncSession,
    *,
    collection: KnowledgeCollection,
    title: str,
    source_kind: DocSourceKind,
    source_uri: str | None,
    raw_text: str | None,
    metadata_json: dict,
    created_by: uuid.UUID | None,
) -> IngestResult:
    """Create the doc row, run the pipeline inline, and commit chunks.

    For the MVP we do this synchronously inside the request. Real deployments
    should move this to the arq worker — the code path is identical.
    """
    # 1) Create a PENDING doc row so the UI can show it immediately.
    repo = AsyncRepository(session, KnowledgeDoc)
    doc = await repo.create(
        collection_id=collection.id,
        title=title,
        source_kind=source_kind,
        source_uri=source_uri,
        raw_text=raw_text,
        status=DocStatus.INGESTING,
        chunk_count=0,
        metadata_json=metadata_json,
        created_by=created_by,
    )

    try:
        if source_kind == DocSourceKind.URL and source_uri:
            raw_text = await _fetch_url(source_uri)
            doc.raw_text = raw_text
        if not raw_text or not raw_text.strip():
            doc.status = DocStatus.FAILED
            doc.error = "empty_text"
            await session.flush([doc])
            return IngestResult(doc=doc, chunks=0)

        config = collection.config_json or {}
        size = int(config.get("chunk_size") or DEFAULT_CHUNK_SIZE)
        overlap = int(config.get("chunk_overlap") or DEFAULT_CHUNK_OVERLAP)
        pieces = _chunk(raw_text, size=size, overlap=overlap)

        total_chars = 0
        for ord_, piece in enumerate(pieces):
            vec, model_tag = await embed_text(piece, workspace_id=collection.workspace_id)
            chunk = KnowledgeChunk(
                doc_id=doc.id,
                collection_id=collection.id,
                ord=ord_,
                text=piece,
                embedding=vec,
                embed_model=model_tag or None,
                tokens=len(piece) // 4,  # rough char→token
            )
            session.add(chunk)
            total_chars += len(piece)

        doc.chunk_count = len(pieces)
        doc.status = DocStatus.READY
        doc.error = None
        await session.flush([doc])
        await session.refresh(doc)  # materialize updated_at before response
        log.info(
            "ingested doc %s chunks=%d chars=%d model=%s",
            doc.id,
            len(pieces),
            total_chars,
            model_tag if pieces else "n/a",
        )
        return IngestResult(doc=doc, chunks=len(pieces))
    except Exception as e:
        log.exception("ingestion failed")
        doc.status = DocStatus.FAILED
        doc.error = str(e)[:2000]
        await session.flush([doc])
        await session.refresh(doc)
        return IngestResult(doc=doc, chunks=0)


# ─── Attachment → RAG bridge ──────────────────────────────
# Mime types we can coerce to UTF-8 text without any external lib. Anything
# outside of this list that is NOT ``application/pdf`` gets a friendly
# ``unsupported_mime`` error so the UI can surface it cleanly.
_TEXTUAL_MIME_PREFIXES: tuple[str, ...] = ("text/",)
_TEXTUAL_MIME_EXACT: frozenset[str] = frozenset(
    [
        "application/json",
        "application/yaml",
        "application/x-yaml",
        "application/xml",
        "application/javascript",
        "application/x-sh",
        "application/x-python-code",
        "application/sql",
        "application/toml",
        "application/x-toml",
    ]
)

# Max bytes we're willing to shove into the chunker in one go. 15 MB covers
# the vast majority of real-world office docs (decks with embedded images,
# multi-sheet spreadsheets) while staying under the 25 MB upload cap.
MAX_EXTRACT_BYTES = 15 * 1024 * 1024


class AttachmentExtractError(ValueError):
    """Raised when an attachment can't be turned into RAG text."""

    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code


def _is_textual_mime(mime: str) -> bool:
    lower = (mime or "").lower()
    if any(lower.startswith(p) for p in _TEXTUAL_MIME_PREFIXES):
        return True
    return lower in _TEXTUAL_MIME_EXACT


# Filename-extension → mime fallback. Browsers / OSes frequently upload office
# documents as ``application/octet-stream`` (or omit the type entirely), and
# Linux's ``mimetypes`` table often lacks the OOXML types. Without this
# fallback those files never reach the right extractor.
_EXT_TO_MIME: dict[str, str] = {
    "pdf": "application/pdf",
    "doc": "application/msword",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xls": "application/vnd.ms-excel",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "ppt": "application/vnd.ms-powerpoint",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv": "text/csv",
    "json": "application/json",
    "yaml": "application/yaml",
    "yml": "application/yaml",
    "xml": "application/xml",
    "md": "text/markdown",
    "markdown": "text/markdown",
    "txt": "text/plain",
}


def _effective_mime(att: Attachment) -> str:
    """Resolve the mime to dispatch on, falling back to the filename extension.

    A trustworthy mime wins; a missing / generic one (``octet-stream`` / ``zip``,
    the latter because OOXML files are zip containers) is overridden by the
    extension when we recognize it.
    """
    mime = (att.mime_type or "").lower()
    if mime and mime not in ("application/octet-stream", "application/zip"):
        return mime
    name = (att.filename or "").lower()
    dot = name.rfind(".")
    if 0 <= dot < len(name) - 1:
        ext = name[dot + 1 :]
        if ext in _EXT_TO_MIME:
            return _EXT_TO_MIME[ext]
    return mime


def extract_text_from_attachment(att: Attachment, data: bytes) -> str:
    """Best-effort text extraction driven by the attachment's mime type.

    Supports:
      * ``text/*`` and common code / config mime types → UTF-8 decode (with
        ``errors="replace"`` so garbage bytes surface as ``?`` rather than a
        hard failure).
      * ``application/pdf`` → pypdf (lazy import; if the library isn't
        installed we raise ``unsupported_mime`` so the UI can show a clear
        "install pypdf" hint).

    Audio / video / image are rejected here — they need transcription / OCR
    which is a separate capability (``fu-image-gen``).
    """
    if att.kind in (AttachmentKind.AUDIO, AttachmentKind.VIDEO, AttachmentKind.IMAGE):
        raise AttachmentExtractError(
            "unsupported_kind",
            f"cannot extract text from attachment kind={att.kind.value}",
        )
    if len(data) > MAX_EXTRACT_BYTES:
        raise AttachmentExtractError(
            "file_too_large",
            f"file > {MAX_EXTRACT_BYTES // (1024 * 1024)}MB; split or compress first",
        )

    mime = _effective_mime(att)

    if _is_textual_mime(mime):
        try:
            return data.decode("utf-8", errors="replace")
        except Exception as e:  # pragma: no cover
            raise AttachmentExtractError("decode_failed", str(e)) from e

    if mime == "application/pdf":
        try:
            from pypdf import PdfReader
        except ImportError as e:
            raise AttachmentExtractError(
                "pdf_lib_missing",
                "pypdf not installed; add it to backend deps to ingest PDFs",
            ) from e
        try:
            reader = PdfReader(io.BytesIO(data))
        except Exception as e:
            raise AttachmentExtractError("pdf_parse_failed", str(e)) from e
        pages: list[str] = []
        for i, page in enumerate(reader.pages):
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                # Don't fail the whole doc on one bad page.
                log.warning("pdf page %d extraction failed for att=%s", i, att.id)
                pages.append("")
        out = "\n\n".join(p for p in pages if p.strip())
        if not out.strip():
            raise AttachmentExtractError(
                "pdf_empty", "pdf produced no extractable text (scanned image?)"
            )
        return out

    if mime in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }:
        return _extract_docx(att, data)

    if mime in {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    }:
        return _extract_xlsx(att, data)

    if mime in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }:
        return _extract_pptx(att, data)

    raise AttachmentExtractError(
        "unsupported_mime",
        f"mime type {mime!r} can't be ingested as text",
    )


def _extract_docx(att: Attachment, data: bytes) -> str:
    """Extract paragraphs + table cells from a .docx file via python-docx.

    ``application/msword`` (legacy .doc binary format) silently falls through
    to python-docx; if the bytes are actually an old binary .doc the
    library raises and we surface ``docx_parse_failed`` with a hint.
    """
    try:
        from docx import Document
    except ImportError as e:
        raise AttachmentExtractError(
            "docx_lib_missing",
            "python-docx not installed; add it to backend deps to ingest DOCX",
        ) from e
    try:
        doc = Document(io.BytesIO(data))
    except Exception as e:
        raise AttachmentExtractError(
            "docx_parse_failed",
            f"docx parse failed (legacy .doc binaries aren't supported — re-save as .docx): {e}",
        ) from e

    parts: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                parts.append(row_text)
    out = "\n".join(parts)
    if not out.strip():
        raise AttachmentExtractError(
            "docx_empty",
            "docx produced no extractable text (image-only document?)",
        )
    return out


def _extract_pptx(att: Attachment, data: bytes) -> str:
    """Extract slide text (shapes + tables) from a .pptx deck via python-pptx.

    Legacy binary ``.ppt`` files aren't supported by python-pptx and surface
    ``pptx_parse_failed`` with a hint to re-save as .pptx.
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise AttachmentExtractError(
            "pptx_lib_missing",
            "python-pptx not installed; add it to backend deps to ingest PPTX",
        ) from e
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise AttachmentExtractError(
            "pptx_parse_failed",
            f"pptx parse failed (legacy .ppt binaries aren't supported — re-save as .pptx): {e}",
        ) from e

    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = (shape.text_frame.text or "").strip()
                if text:
                    slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        slide_parts.append(row_text)
        if slide_parts:
            parts.append(f"# Slide {idx}\n" + "\n".join(slide_parts))
    out = "\n\n".join(parts)
    if not out.strip():
        raise AttachmentExtractError(
            "pptx_empty",
            "pptx produced no extractable text (image-only deck?)",
        )
    return out


def _extract_xlsx(att: Attachment, data: bytes) -> str:
    """Extract every sheet as a tab-separated table via openpyxl.

    Empty rows / cells are skipped; each sheet header line records the
    sheet name so the model can reference rows in context.
    """
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise AttachmentExtractError(
            "xlsx_lib_missing",
            "openpyxl not installed; add it to backend deps to ingest XLSX",
        ) from e
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:
        raise AttachmentExtractError("xlsx_parse_failed", str(e)) from e

    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
        parts.append("")
    out = "\n".join(parts).strip()
    if not out:
        raise AttachmentExtractError(
            "xlsx_empty", "xlsx produced no extractable text (empty workbook?)"
        )
    return out


async def ingest_attachment(
    session: AsyncSession,
    *,
    collection: KnowledgeCollection,
    attachment: Attachment,
    data: bytes,
    created_by: uuid.UUID | None,
    title_override: str | None = None,
) -> IngestResult:
    """Extract text from an attachment and feed it into ``ingest_document``.

    Stores the attachment id + filename in ``metadata_json`` so the knowledge
    page can link back to the original blob.
    """
    raw_text = extract_text_from_attachment(attachment, data)
    title = (title_override or attachment.filename or "attachment")[:255]
    return await ingest_document(
        session,
        collection=collection,
        title=title,
        source_kind=DocSourceKind.FILE,
        source_uri=f"attachment://{attachment.id}",
        raw_text=raw_text,
        metadata_json={
            "attachment_id": str(attachment.id),
            "filename": attachment.filename,
            "mime_type": attachment.mime_type,
            "size_bytes": attachment.size_bytes,
            "sha256": attachment.sha256,
        },
        created_by=created_by,
    )


# ─── Search ────────────────────────────────────────────────
@dataclass
class ChunkHit:
    id: uuid.UUID
    doc_id: uuid.UUID
    doc_title: str | None
    ord: int
    text: str
    score: float


async def search(
    session: AsyncSession,
    *,
    collection: KnowledgeCollection,
    query: str,
    top_k: int = 5,
    allowed_doc_ids: set[uuid.UUID] | None = None,
) -> list[ChunkHit]:
    """Semantic search over a collection.

    ``allowed_doc_ids`` is an optional document-level ACL filter (resolved by
    :func:`app.services.kb_source.filter_accessible_doc_ids`). ``None`` means
    no filtering; an empty set returns no results — the caller has no access
    to any doc in the collection.
    """
    vec, _ = await embed_text(query, workspace_id=collection.workspace_id)
    if allowed_doc_ids is not None and not allowed_doc_ids:
        return []
    if vec is None:
        # No embedder configured for this workspace — semantic search is
        # disabled (no silent hash fallback). Callers see an empty
        # result set and can fall back to a different surface.
        return []
    # pgvector's ``<=>`` is cosine *distance* (smaller = better). Convert to
    # similarity = 1 - distance so the UI sees an intuitive 0-1 score.
    distance = KnowledgeChunk.embedding.cosine_distance(vec)
    stmt = (
        select(
            KnowledgeChunk.id,
            KnowledgeChunk.doc_id,
            KnowledgeDoc.title,
            KnowledgeChunk.ord,
            KnowledgeChunk.text,
            distance.label("distance"),
        )
        .join(KnowledgeDoc, KnowledgeDoc.id == KnowledgeChunk.doc_id)
        .where(
            KnowledgeChunk.collection_id == collection.id,
            KnowledgeChunk.embedding.is_not(None),
            KnowledgeDoc.deleted_at.is_(None),
        )
        .order_by(distance.asc())
        .limit(top_k)
    )
    if allowed_doc_ids is not None:
        stmt = stmt.where(KnowledgeChunk.doc_id.in_(allowed_doc_ids))
    rows = (await session.execute(stmt)).all()
    return [
        ChunkHit(
            id=r.id,
            doc_id=r.doc_id,
            doc_title=r.title,
            ord=r.ord,
            text=r.text,
            score=float(1 - (r.distance or 1.0)),
        )
        for r in rows
    ]


# ─── Helpers ──────────────────────────────────────────────
async def _fetch_url(url: str) -> str:
    """Download URL → readable text via trafilatura."""
    try:
        import trafilatura
    except ImportError:
        trafilatura = None  # type: ignore[assignment]

    async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as c:
        r = await c.get(url, headers={"User-Agent": "SenHarness/0.1 (+rag)"})
        r.raise_for_status()
        html = r.text

    if trafilatura is not None:
        extracted = trafilatura.extract(
            html, include_comments=False, include_tables=True, output_format="markdown"
        )
        if extracted and extracted.strip():
            return extracted
    # Fallback: strip tags crudely.
    import re

    no_scripts = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", html)
    txt = re.sub(r"(?s)<[^>]+>", " ", no_scripts)
    return re.sub(r"\s+", " ", txt).strip()


def _chunk(text: str, *, size: int, overlap: int) -> list[str]:
    """Naive char-window chunker. Fine for MVP.

    For better quality later, swap in a paragraph+sentence-aware chunker
    (``pydantic-ai-knowledge`` ships one).
    """
    text = (text or "").strip()
    if not text:
        return []
    if len(text) <= size:
        return [text]

    step = max(size - overlap, 1)
    out: list[str] = []
    i = 0
    while i < len(text):
        piece = text[i : i + size].strip()
        if piece:
            out.append(piece)
        i += step
    return out
